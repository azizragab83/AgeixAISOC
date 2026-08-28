# -*- coding: utf-8 -*-
"""Fill the official Final Discussion PPT template + Poster template for AgeixAISOC."""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

DOWNLOADS = r"C:\Users\Digilians\Downloads"
DOC_DIR = os.path.join(DOWNLOADS, "doc")
BUILD_DIR = os.path.dirname(os.path.abspath(__file__))
os.makedirs(DOC_DIR, exist_ok=True)

CYAN = RGBColor(0x06, 0xB6, 0xD4)
DARK = RGBColor(0x1E, 0x29, 0x3B)

TEAM = "Aziz Ragab (Team Leader)   |   Mohamed Hany   |   Emad Hassan   |   Taha Elghonaimi"
SUPERVISOR = "General Supervisor: Dr. Rabab M. Nabawy        Academic Director: Dr. Ahmed Tobal"


def set_body(tf, lines, size=18, color=DARK, bold_first=False):
    tf.clear()
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        if line and not bold_first and not line.startswith(" ") and not line.endswith(":"):
            run.text = "- " + line
        else:
            run.text = line
        f = run.font
        f.size = Pt(size)
        f.name = "Calibri"
        f.color.rgb = color
        para.space_after = Pt(4)


def find_body(slide):
    fallback = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.is_placeholder and sh.placeholder_format.idx != 0 and "title" not in sh.name.lower():
            return sh
        if fallback is None and not sh.is_placeholder:
            fallback = sh
    return fallback


def fill(slide, lines, size=18, bold_first=False):
    sh = find_body(slide)
    if sh is not None:
        set_body(sh.text_frame, lines, size=size, bold_first=bold_first)
    return sh


def build_presentation():
    prs = Presentation(os.path.join(DOWNLOADS, "Final Discussion_Cybersecurity.pptx"))
    slides = list(prs.slides)

    # Slide 1 - title boxes
    s = slides[0]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if "Project Title" in t or "9-Month" in t:
            set_body(sh.text_frame, [
                "Digilians 9-Month Diploma",
                "AgeixAISOC",
                "AI-Orchestrated Cyber Defense Platform",
                "Track: Cybersecurity",
                SUPERVISOR,
            ], size=20)
            try:
                r = sh.text_frame.paragraphs[1].runs[0]
                r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = CYAN
            except Exception:
                pass
        elif "Presented by" in t:
            set_body(sh.text_frame, ["Presented by:", TEAM], size=18)

    fill(slides[1], [
        "Introduction - why AI belongs in the SOC",
        "Problem Definition - alert fatigue, fragmentation, autonomy risk",
        "Motivation & Objectives",
        "Implementation - stack, lab, agents",
        "Architecture - end-to-end closed loop",
        "Methodology - how we built & tested it",
        "Results - what works today",
        "Future Work  |  References",
    ], size=22)

    fill(slides[2], [
        "Modern SOCs cannot scale manually against growing alert volume.",
        "AI helps, but fully autonomous defense creates legal accountability gaps.",
        "AgeixAISOC: an AI-orchestrated SOC platform running 100% locally (data sovereignty).",
        "Core principle:",
        "\"AI orchestrates . AI detects . AI recommends -- HUMAN DECIDES -- System executes.\"",
        "Built with FastAPI, LangGraph, CrewAI, Ollama local LLMs, ChromaDB RAG, React dashboard.",
    ], size=20)

    fill(slides[3], [
        "Alert Fatigue - roughly 70% of SIEM alerts are false positives.",
        "Tool Fragmentation - analysts pivot across SIEM / EDR / NDR / TI consoles.",
        "Slow Response - industry MTTR measured in hours; attackers move freely.",
        "Autonomy Risk - self-acting AI lacks human accountability and legal defensibility.",
        "Detection Gaps - blind spots usually found only during periodic manual pentesting.",
    ], size=20)

    fill(slides[4], [
        "Motivation:",
        "Prove that a coordinated multi-agent AI pipeline can triage, score and recommend while the human keeps final authority.",
        "Keep every byte local (Ollama) - data sovereignty for regulated environments.",
        "Objectives:",
        "Ingest live Wazuh SIEM alerts through a custom integration webhook.",
        "Analyze with a LangGraph-orchestrated agent pipeline producing one unified decision package.",
        "Enforce Human-in-the-Loop approval before ANY action.",
        "Execute approved blocks on FortiGate via n8n SOAR (with direct-API fallback).",
        "Learn from every analyst decision via adaptive RAG memory (ChromaDB).",
    ], size=19, bold_first=True)

    fill(slides[5], [
        "Backend - Python 3.11 / FastAPI: 18 endpoints, WebSocket broadcast, dual-import support.",
        "Orchestrator - LangGraph StateGraph; 7 CrewAI agents (detection, risk scoring, recommendation, threat hunter, forensics, red team, detection engineering) with 30s timeouts and graceful fallbacks.",
        "Memory - ChromaDB RAG stores every human decision; similar future alerts auto-suppressed or down-scored.",
        "Live Lab - VMware 192.168.56.0/24: Kali attacker (SSH bridge), Wazuh SIEM + custom webhook, Windows 10 / AD victims, Metasploitable2.",
        "SOAR - n8n workflow executes FortiGate REST blocks; FastAPI falls back to direct FortiGate calls if n8n is down.",
        "Frontend - React/Vite tactical dashboard: live terminal, HITL cards, MITRE coverage, AR/EN bilingual, dark/light themes.",
    ], size=17)

    # Slide 7 architecture: image if generated, else text
    s7 = slides[6]
    arch_img = os.path.join(BUILD_DIR, "architecture.png")
    from pptx.util import Inches
    body = None
    for sh in s7.shapes:
        if sh.has_text_frame and sh.is_placeholder and sh.placeholder_format.idx != 0:
            body = sh
            break
    if os.path.exists(arch_img):
        if body is not None and body.has_text_frame and not body.text_frame.text.strip():
            s7.shapes.add_picture(arch_img, Inches(1.0), Inches(2.0), width=Inches(11.8))
        elif body is not None:
            # keep text but shrink, add image below title area anyway
            set_body(body.text_frame, [
                "Kali attack -> Wazuh SIEM alert (webhook) -> Agent Pipeline (LangGraph)",
                "-> Cognitive synthesis -> Decision Package -> HITL Dashboard (Approve/Reject)",
                "-> n8n SOAR -> FortiGate block -> Audit + RAG learning feedback",
            ], size=16)
    else:
        fill(s7, [
            "Kali attack -> Wazuh SIEM alert (custom webhook)",
            "-> LangGraph agent pipeline (7 agents, cognitive synthesis)",
            "-> Decision Package -> HITL Dashboard (Approve / Reject)",
            "-> n8n SOAR -> FortiGate REST block",
            "-> Audit trail + adaptive RAG learning feedback",
        ], size=22)

    fill(slides[7], [
        "Iterative, test-driven development inside an isolated cyber range (Host-Only 192.168.56.0/24).",
        "Real attacks: Nmap scans and Hydra brute-force launched from Kali over SSH by the platform itself.",
        "Wazuh custom integration forwards every alert at level 7+ as JSON to the backend webhook.",
        "Every agent step is time-boxed (30s), returns structured JSON, degrades gracefully - pipeline never crashes.",
        "HITL governance: approve/reject recorded in audit trail; every decision trains the adaptive RAG.",
        "Validation: 107 automated pytest tests plus scripted live end-to-end scenarios.",
    ], size=19)

    fill(slides[8], [
        "End-to-end pipeline validated LIVE: real attack -> SIEM alert -> agent analysis -> human approval -> FortiGate block.",
        "Test suite: 107 automated tests, 104 passing; the 3 failures are environment-only SSH timeouts (VMs offline during CI).",
        "Adaptive learning verified: rejected alerts stored in RAG; similar future alerts auto-suppressed (>90% match) or down-scored (70%+).",
        "Zero-downtime SOAR: n8n failure triggers automatic direct FortiGate fallback - response never stops.",
        "Lab-scale performance: detection-to-decision under one minute; targets MTTD < 5 min, MTTR < 20 min.",
        "Honest scope: UEBA and deception engines exist as wired modules/endpoints - full autonomous operation is declared future work.",
    ], size=18)

    fill(slides[9], [
        "Promote UEBA (Isolation Forest) and deception/honeypot redirection from wired endpoints to fully autonomous operation.",
        "Register our fine-tuned local model (ageix-brain, Unsloth QLoRA to GGUF) as Master Brain default.",
        "Expand threat-intel connectors: AlienVault OTX, AbuseIPDB, MISP feeds.",
        "Multi-tenant SaaS hardening: JWT auth, RBAC roles, PostgreSQL row-level security.",
        "Sigma-rule auto-deployment loop: gap detection -> rule generation -> human approval -> Wazuh push.",
        "Cyber digital-twin simulation before executing high-impact responses.",
    ], size=19)

    fill(slides[10], [
        "[1] MITRE ATT&CK Framework - https://attack.mitre.org",
        "[2] Wazuh Open Source XDR Documentation - https://documentation.wazuh.com",
        "[3] CrewAI Multi-Agent Framework - https://github.com/crewAIInc/crewAI",
        "[4] LangGraph Stateful Agent Workflows - https://langchain-ai.github.io/langgraph",
        "[5] Ollama Local LLM Runtime - https://ollama.com",
        "[6] ChromaDB Embedding Database - https://www.trychroma.com",
        "[7] n8n Workflow Automation Docs - https://docs.n8n.io",
        "[8] FortiOS REST API Reference - https://docs.fortinet.com",
        "[9] FastAPI Documentation - https://fastapi.tiangolo.com",
        "[10] Sigma Generic Signature Format - https://github.com/SigmaHQ/sigma",
    ], size=16)

    sh12 = find_body(slides[11])
    if sh12 is not None:
        set_body(sh12.text_frame, ["Thank You", "", "Questions?", "", TEAM], size=28)
        try:
            r = sh12.text_frame.paragraphs[0].runs[0]
            r.font.bold = True; r.font.size = Pt(48); r.font.color.rgb = CYAN
        except Exception:
            pass

    out = os.path.join(DOC_DIR, "Final_Presentation_AgeixAISOC.pptx")
    prs.save(out)
    print("[OK]", out)


if __name__ == "__main__":
    build_presentation()
