# -*- coding: utf-8 -*-
"""Fill the official Grad Poster template with AgeixAISOC content."""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

DOWNLOADS = r"C:\Users\Digilians\Downloads"
DOC_DIR = os.path.join(DOWNLOADS, "doc")
BUILD_DIR = os.path.dirname(os.path.abspath(__file__))

DARK = RGBColor(0x1E, 0x29, 0x3B)
CYAN = RGBColor(0x0E, 0x74, 0x90)

TEAM = ("Aziz Ragab (Team Leader)   Mohamed Hany   Emad Hassan   Taha Elghonaimi")

ABSTRACT = (
    "Security Operations Centers face severe alert fatigue (about 70% false positives), fragmented tooling, "
    "and legal risks from fully autonomous AI. AgeixAISOC is an enterprise-grade, AI-orchestrated cyber defense "
    "platform that runs 100% locally to guarantee data sovereignty. A LangGraph-orchestrated swarm of specialized "
    "CrewAI agents - each augmented with adaptive RAG memory and live web search - analyzes live Wazuh SIEM "
    "telemetry, performs cognitive synthesis, and issues a single unified decision package. Strict Human-in-the-Loop "
    "(HITL) governance ensures no action executes without analyst approval; approved blocks are executed on FortiGate "
    "through an n8n SOAR engine with a zero-downtime direct-API fallback. Every human decision feeds the RAG memory, "
    "progressively eliminating false positives."
)

AIM = (
    "Eliminate alert fatigue and accelerate incident response while preserving 100% human accountability: "
    "detect -> analyze -> recommend -> human decides -> system executes. Target metrics: MTTD under 5 minutes, "
    "MTTR under 20 minutes, zero unauthorized actions."
)

DATASET = (
    "Live VMware cyber range (192.168.56.0/24):\n"
    "- Kali Linux: real attacks (Nmap, Hydra) launched via SSH bridge\n"
    "- Wazuh SIEM: Sysmon/auth telemetry; custom webhook forwards alerts level 7+\n"
    "- Windows 10 + Active Directory + Metasploitable2 victim assets\n"
    "- MITRE ATT&CK technique mapping for coverage analysis"
)

METHODOLOGY = (
    "1. Attack simulation - backend triggers real attacks from Kali over SSH.\n"
    "2. Detection - Wazuh fires alerts; custom integration POSTs JSON to FastAPI.\n"
    "3. Agent pipeline - LangGraph runs 7 time-boxed CrewAI agents (detection, risk scoring, recommendation, threat hunter, forensics, red team, detection engineering).\n"
    "4. Cognitive synthesis - Master Brain merges agent outputs with RAG recall and live web search into one decision package.\n"
    "5. HITL - analyst Approves or Rejects on the dashboard; every decision is audited.\n"
    "6. SOAR - n8n executes the FortiGate REST block; direct-API fallback guarantees uptime.\n"
    "7. Learning - decisions embedded into ChromaDB; similar alerts later suppressed automatically."
)

RESULTS = (
    "- Full closed loop validated live: attack -> SIEM -> AI -> human -> firewall block.\n"
    "- 107 automated tests (104 passing; 3 failures are VM-offline SSH timeouts only).\n"
    "- Adaptive RAG proven: rejected patterns auto-suppressed on recurrence (over 90% similarity).\n"
    "- Zero-downtime SOAR fallback verified (n8n outage triggers direct FortiGate call).\n"
    "- Detection-to-decision latency under one minute at lab scale."
)

DEMO = (
    "Press Launch Attack on the dashboard -> Kali runs a real Nmap/Hydra attack -> Wazuh alerts -> "
    "agents stream their reasoning to the live terminal -> HITL card appears -> analyst clicks Approve -> "
    "FortiGate blocks the attacker IP. All under human control."
)

CONCLUSION = (
    "AgeixAISOC transforms the SOC from a reactive alert factory into a proactive, self-correcting defense system: "
    "multi-agent AI performs the heavy analysis, the human stays accountable, and the platform executes - with every "
    "step audited and learned from."
)

FUTURE = (
    "- UEBA anomaly engine (Isolation Forest) in fully autonomous operation\n"
    "- Deception / honeypot redirection (Ghost VLAN)\n"
    "- Fine-tuned local model (Unsloth QLoRA) as Master Brain default\n"
    "- Multi-tenant SaaS hardening: JWT auth, RBAC, row-level security\n"
    "- Sigma-rule gap loop: detect gap -> generate rule -> approve -> deploy to Wazuh\n"
    "- Cyber digital-twin simulation before high-impact responses"
)

TOOLS = (
    "Python 3.11 | FastAPI | LangGraph | CrewAI | Ollama (qwen2.5:14b, llama3.1, qwen2.5-coder) | "
    "ChromaDB | Wazuh SIEM + custom integration webhook | FortiGate REST API | n8n SOAR | React + Vite + Tailwind | "
    "WebSockets | Paramiko SSH | pytest (107 tests)"
)


def set_text(shape, text, size=22):
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        para.space_after = Pt(6)


def build_poster():
    prs = Presentation(os.path.join(DOWNLOADS, "Grad_Poster_Cybersecurity.pptx"))
    slide = prs.slides[0]

    # map shapes by name
    by_name = {}
    for sh in slide.shapes:
        key = sh.name
        if key not in by_name:
            by_name[key] = sh
        else:
            # duplicate names: keep both in list
            by_name[key] = [by_name[key], sh] if not isinstance(by_name[key], list) else by_name[key] + [sh]

    def get(name, idx=0):
        v = by_name.get(name)
        if isinstance(v, list):
            return v[idx] if idx < len(v) else None
        return v

    # Title & team
    t = get("Google Shape;85;p1")
    if t is not None:
        set_text(t, "AgeixAISOC", size=44)
        p2 = t.text_frame.add_paragraph()
        r = p2.add_run(); r.text = "AI-Orchestrated Cyber Defense Platform"
        r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = CYAN

    pb = get("Google Shape;86;p1", 0)
    if pb is not None and pb.has_text_frame and "Presented" in pb.text_frame.text:
        set_text(pb, "Presented by:\n" + TEAM, size=16)

    psup = get("Google Shape;86;p1", 1)
    if psup is not None:
        set_text(psup, "Under Supervision of:\nDr. Rabab M. Nabawy (General Supervisor)\nDr. Ahmed Tobal (Academic Director)", size=14)

    # Abstract body
    ab = get("TextBox 1")
    if ab is not None:
        set_text(ab, ABSTRACT, size=18)

    # Aim
    aim = get("TextBox 10")
    if aim is not None:
        set_text(aim, AIM, size=20)

    # Dataset
    ds = get("TextBox 2")
    if ds is not None:
        set_text(ds, DATASET, size=15)

    # Methodology
    mt = get("TextBox 12")
    if mt is not None:
        set_text(mt, METHODOLOGY, size=13)

    # Results
    rs = get("TextBox 19")
    if rs is not None:
        set_text(rs, RESULTS, size=14)

    # Demo flow (body box under 'Demos' header)
    demo_box = get("Google Shape;105;p1")
    if demo_box is not None and demo_box.has_text_frame:
        set_text(demo_box, DEMO, size=14)

    # Conclusion
    cc = get("TextBox 22")
    if cc is not None:
        set_text(cc, CONCLUSION, size=14)

    # Future work box (the one containing placeholder sentence)
    fw = get("Google Shape;108;p1", 0)
    if fw is not None and fw.has_text_frame and "dataset size" in fw.text_frame.text.lower():
        set_text(fw, FUTURE, size=13)

    # Used tools empty box (duplicate name at tools area)
    tools = get("Google Shape;108;p1", 1)
    if tools is not None and tools.has_text_frame and not tools.text_frame.text.strip():
        set_text(tools, TOOLS, size=12)
    elif tools is None:
        # fallback: find any empty shape near bottom-left
        for sh in slide.shapes:
            if sh.has_text_frame and not sh.text_frame.text.strip() and sh.top > Emu(int(110 * 360000)) and sh.left < Emu(int(30 * 360000)):
                set_text(sh, TOOLS, size=12)
                break

    # Remove glaucoma screenshots, insert our architecture diagram in SYSTEM OVERVIEW zone
    arch_img = os.path.join(BUILD_DIR, "architecture.png")
    overview_pic = get("Picture 3")
    if overview_pic is not None:
        left, top, width, height = overview_pic.left, overview_pic.top, overview_pic.width, overview_pic.height
        overview_pic._element.getparent().remove(overview_pic._element)
        if os.path.exists(arch_img):
            slide.shapes.add_picture(arch_img, left, top, width=width, height=height)

    # Remove leftover glaucoma pictures
    for name in ["Picture 11", "Picture 14", "Picture 16", "Picture 21"]:
        pic = get(name)
        if pic is not None:
            try:
                pic._element.getparent().remove(pic._element)
            except Exception:
                pass

    out = os.path.join(DOC_DIR, "Grad_Poster_AgeixAISOC.pptx")
    prs.save(out)
    print("[OK]", out)


if __name__ == "__main__":
    build_poster()
